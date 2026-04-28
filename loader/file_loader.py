import io
import re
from PIL import Image

from docx import Document
from pptx import Presentation
from pypdf import PdfReader

from loader.ocr import ocr_image_bytes


# ==============================
# Helpers
# ==============================
def _clean(text: str) -> str:
    text = text.replace("\n", " ")
    text = re.sub(r"\s+", " ", text)
    return text.strip()


# ==============================
# MAIN EXTRACT
# ==============================
def extract_text(file):
    name = file.name.lower()
    raw = file.getvalue()

    texts = []
    ocr_texts = []

    # ==================================================
    # TXT
    # ==================================================
    if name.endswith(".txt"):
        texts.append(raw.decode("utf-8", errors="ignore"))

    # ==================================================
    # IMAGE FILES → OCR 100%
    # ==================================================
    elif name.endswith((".png", ".jpg", ".jpeg", ".webp")):
        ocr_texts.append(ocr_image_bytes(raw))

    # ==================================================
    # PDF
    # ==================================================
    elif name.endswith(".pdf"):
        reader = PdfReader(io.BytesIO(raw))

        for page in reader.pages:
            # 1) text layer
            t = page.extract_text()
            if t:
                texts.append(t)

        # 2) OCR toàn bộ PDF (scan + image)
        try:
            ocr_texts.append(ocr_image_bytes(raw))
        except Exception as e:
            print("PDF OCR error:", e)

    # ==================================================
    # DOCX
    # ==================================================
    elif name.endswith(".docx"):
        doc = Document(io.BytesIO(raw))

        # text
        for p in doc.paragraphs:
            if p.text.strip():
                texts.append(p.text)

        # OCR images inside DOCX
        for rel in doc.part._rels:
            rel = doc.part._rels[rel]
            if "image" in rel.target_ref:
                img_bytes = rel.target_part.blob
                ocr_texts.append(ocr_image_bytes(img_bytes))

    # ==================================================
    # PPTX
    # ==================================================
    elif name.endswith(".pptx"):
        prs = Presentation(io.BytesIO(raw))
        for slide in prs.slides:
            for shape in slide.shapes:
                # text
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)

                # image → OCR
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    img_bytes = shape.image.blob
                    ocr_texts.append(ocr_image_bytes(img_bytes))

    else:
        raise ValueError(f"Unsupported file type: {name}")

    # ==================================================
    # MERGE TEXT + OCR
    # ==================================================
    final_text = "\n\n".join(texts + ocr_texts)
    final_text = _clean(final_text)

    return final_text